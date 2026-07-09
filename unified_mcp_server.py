from __future__ import annotations  # 型注釈を遅延評価に（tree-sitter未導入でも import が通るように）

import json
import sys
import hashlib
import re
import os
import unicodedata
import threading
import subprocess
import shlex
from pathlib import Path
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed

from mcp.server.fastmcp import FastMCP

# PyMuPDF（PDF抽出用 — 未インストールでもサーバーは起動し、PDF機能だけ無効化される）
try:
    import fitz  # PyMuPDF
    _PYMUPDF_AVAILABLE = True
except ImportError:
    fitz = None
    _PYMUPDF_AVAILABLE = False

# chromadb / sentence-transformers は RAG 有効時のみ遅延importする（重い依存）
# → _get_collection() の中で読み込む

# tree-sitter (symbol解析用 — 未インストールでもサーバーは起動する)
try:
    from tree_sitter import Language, Parser as TSParser
    _TREE_SITTER_AVAILABLE = True
except ImportError:
    _TREE_SITTER_AVAILABLE = False

# ========================================================
# 設定
# ========================================================

# --- すべて環境変数 / workspace_config.json で上書き可能（ハードコードなし） ---
#
#   WORKSPACE_ROOT : 触らせる作業フォルダのルート（既定: このリポジトリの ./workspace）
#     └─ apps/<project>/  … 読み書き・grep・編集の対象
#     └─ docs/<project>/  … RAG（意味検索）の対象
#   CHROMA_DB_DIR  : ベクタDBの保存先（既定: <WORKSPACE_ROOT>/.chroma_db）
#   ENABLE_RAG     : "0" で RAG 系ツールを無効化（重い依存を読み込まない）
#   EMBED_MODEL    : 埋め込みモデル名（既定: intfloat/multilingual-e5-large）
#
_REPO_DIR = Path(__file__).resolve().parent

MCP_ROOT = os.environ.get("WORKSPACE_ROOT") or str(_REPO_DIR / "workspace")
APPS_ROOT = os.path.join(MCP_ROOT, "apps")
DOCS_ROOT = os.path.join(MCP_ROOT, "docs")
DB_DIR    = os.environ.get("CHROMA_DB_DIR") or os.path.join(MCP_ROOT, ".chroma_db")

ENABLE_RAG  = os.environ.get("ENABLE_RAG", "1").lower() not in ("0", "false", "no", "")
EMBED_MODEL = os.environ.get("EMBED_MODEL", "intfloat/multilingual-e5-large")

# 追加設定は workspace_config.json（任意）で与える:
#   {
#     "apps_only_projects": ["teach"],                     # RAG不要（appsのみ）のプロジェクト
#     "extra_rag_paths": {"<project>": ["C:/abs/path", ...]}  # docs/以外もRAG対象に足す
#   }
_CFG_PATH = Path(os.environ.get("WORKSPACE_CONFIG") or (_REPO_DIR / "workspace_config.json"))
_cfg: dict = {}
if _CFG_PATH.exists():
    try:
        _cfg = json.loads(_CFG_PATH.read_text(encoding="utf-8"))
    except Exception as _e:
        print(f"[config] {_CFG_PATH.name} の読み込みに失敗: {_e}", file=sys.stderr)

APPS_ONLY_PROJECTS: set[str] = set(_cfg.get("apps_only_projects", []))  # RAGなしプロジェクト
EXTRA_RAG_PATHS: dict[str, list[str]] = _cfg.get("extra_rag_paths", {})

# 除外フォルダ
EXCLUDE_DIRS = {
    "__pycache__", ".git", "node_modules", ".venv",
    ".chroma_db", ".backups", "no_using", "paper_search.egg-info", "zemi_vector_db"
}

CHUNK_SIZE = 500
CHUNK_OVERLAP = 100

# PDF（学術論文）は段落が長いので大きめのチャンクを使う
PDF_CHUNK_SIZE = 1500
PDF_CHUNK_OVERLAP = 300

SEARCH_N_RESULTS = 30
SEARCH_MAX_FILES = 7

SEARCH_DISTANCE_THRESHOLD = 1.5

BINARY_EXTENSIONS = {".pdf", ".pptx", ".docx", ".xlsx", ".xls"}

# --- [FIX①] read_document / read_file の返却上限 ---
MAX_READ_CHARS = 30000  # 約30,000文字。超過時は先頭+末尾を返す

# --- [FIX⑥] search_in_files のテキスト検索対象拡張子ホワイトリスト ---
TEXT_SEARCH_EXTENSIONS = {
    ".py", ".pyi", ".js", ".jsx", ".ts", ".tsx", ".vue", ".svelte",
    ".html", ".htm", ".css", ".scss", ".sass", ".less",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".env",
    ".md", ".mdx", ".txt", ".rst", ".tex", ".csv", ".tsv",
    ".sh", ".bash", ".zsh", ".bat", ".ps1", ".cmd",
    ".sql", ".r", ".R", ".jl", ".lua", ".rb", ".go", ".rs",
    ".java", ".kt", ".kts", ".scala", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".swift", ".m", ".mm", ".dart", ".php",
    ".xml", ".svg", ".graphql", ".gql", ".proto",
    ".dockerfile", ".gitignore", ".editorconfig",
    ".ipynb",
}
# 拡張子なしファイルで名前ベースで許可するもの
TEXT_SEARCH_NAMES = {
    "Makefile", "Dockerfile", "Procfile", "Gemfile", "Rakefile",
    "CMakeLists.txt", ".gitignore", ".env", ".env.local",
}

# --- [FIX⑥] search_in_files のファイルサイズ上限 ---
MAX_SEARCH_FILE_SIZE = 1 * 1024 * 1024  # 1MB

# --- [NEW④] run_check のデフォルトタイムアウト ---
CHECK_TIMEOUT_SECONDS = 60

# --- [NEW④] run_check の設定ファイル名 ---
CHECK_CONFIG_FILENAME = "mcp_checks.json"

# ========================================================
# [NEW①③] Symbol解析 — 言語設定（追加はここに辞書を足すだけ）
# ========================================================

# 各言語の定義:
#   module:       pip install tree-sitter-<name> で入れたモジュール名
#   extensions:   この言語として扱う拡張子
#   definitions:  シンボル定義として抽出するASTノード型
#   name_field:   ノードから名前を取るフィールド名（"name" or "identifier" 等）
#   arrow_detect: arrow function を変数名で拾うかどうか (JS/TS用)
LANGUAGE_CONFIG: dict[str, dict] = {
    "python": {
        "module": "tree_sitter_python",
        "extensions": [".py", ".pyi"],
        "definitions": [
            "function_definition",
            "class_definition",
        ],
        "references": [
            "import_statement",
            "import_from_statement",
        ],
        "name_field": "name",
    },
    "typescript": {
        "module": "tree_sitter_typescript",
        "language_func": "language_typescript",
        "extensions": [".ts"],
        "definitions": [
            "function_declaration",
            "class_declaration",
            "interface_declaration",
            "type_alias_declaration",
            "enum_declaration",
        ],
        "references": [
            "import_statement",
        ],
        "name_field": "name",
        "arrow_detect": True,
    },
    "tsx": {
        "module": "tree_sitter_typescript",
        "language_func": "language_tsx",
        "extensions": [".tsx"],
        "definitions": [
            "function_declaration",
            "class_declaration",
            "interface_declaration",
            "type_alias_declaration",
            "enum_declaration",
        ],
        "references": [
            "import_statement",
        ],
        "name_field": "name",
        "arrow_detect": True,
    },
    "javascript": {
        "module": "tree_sitter_javascript",
        "extensions": [".js", ".jsx", ".mjs"],
        "definitions": [
            "function_declaration",
            "class_declaration",
        ],
        "references": [
            "import_statement",
        ],
        "name_field": "name",
        "arrow_detect": True,
    },
    "matlab": {
        "module": "tree_sitter_matlab",
        "extensions": [".m"],
        "definitions": [
            "function_definition",
        ],
        "references": [],
        "name_field": "name",
    },
}

# --- パーサーの遅延初期化 ---
_parsers: dict = {}
_languages: dict = {}
_ext_to_lang: dict[str, str] = {}  # ".py" -> "python"
_parsers_init_done = False


def _init_parsers():
    """利用可能な言語パーサーを初期化する。未インストールの言語はスキップ。"""
    global _parsers_init_done
    if _parsers_init_done:
        return  # 初期化済み
    _parsers_init_done = True

    if not _TREE_SITTER_AVAILABLE:
        print("[symbol] tree-sitter is not installed. Symbol analysis disabled.", file=sys.stderr)
        return

    import importlib
    for lang_name, config in LANGUAGE_CONFIG.items():
        try:
            mod = importlib.import_module(config["module"])
            lang_func_name = config.get("language_func", "language")
            lang_func = getattr(mod, lang_func_name)
            language = Language(lang_func())
            parser = TSParser(language)

            _parsers[lang_name] = parser
            _languages[lang_name] = language
            for ext in config["extensions"]:
                _ext_to_lang[ext] = lang_name
        except (ImportError, AttributeError, Exception) as e:
            print(f"[symbol] {lang_name} parser unavailable: {e}", file=sys.stderr)


def _require_tree_sitter() -> str | None:
    """tree-sitter が利用可能か確認。利用不可ならエラーメッセージ文字列を返す。"""
    _init_parsers()
    if not _TREE_SITTER_AVAILABLE:
        return (
            "❌ tree-sitter が未インストールです。\n"
            "仮想環境で以下を実行してください:\n\n"
            "  pip install tree-sitter tree-sitter-python tree-sitter-javascript "
            "tree-sitter-typescript tree-sitter-matlab\n\n"
            "インストール後、MCPサーバーを再起動してください。"
        )
    if not _parsers:
        return (
            "❌ 言語パーサーが1つもロードされていません。\n"
            "tree-sitter-python 等の言語パッケージをインストールしてください。"
        )
    return None


def _get_parser_for_file(filepath: Path) -> tuple[str, TSParser, dict] | None:
    """ファイル拡張子からパーサーと言語設定を返す。対応外はNone。"""
    _init_parsers()
    ext = filepath.suffix.lower()
    lang_name = _ext_to_lang.get(ext)
    if lang_name is None:
        return None
    if lang_name not in _parsers:
        return None
    return lang_name, _parsers[lang_name], LANGUAGE_CONFIG[lang_name]


def _extract_symbols_from_file(filepath: Path) -> list[dict] | None:
    """1ファイルからシンボル（関数・クラス等）の定義情報を抽出する。"""
    info = _get_parser_for_file(filepath)
    if info is None:
        return None
    lang_name, parser, config = info

    try:
        source = filepath.read_bytes()
    except Exception:
        return None

    tree = parser.parse(source)
    symbols = []

    def _extract_name(node, config):
        """ノードからシンボル名を取得する"""
        name_field = config.get("name_field", "name")
        name_node = node.child_by_field_name(name_field)
        if name_node:
            return name_node.text.decode("utf-8", errors="replace")
        # フォールバック: 最初の identifier 子ノードを使う
        for child in node.children:
            if child.type == "identifier" or child.type == "type_identifier":
                return child.text.decode("utf-8", errors="replace")
        return None

    def _walk_for_definitions(node, depth=0):
        """ASTを走査してシンボル定義を収集する"""
        if node.type in config["definitions"]:
            name = _extract_name(node, config)
            if name:
                symbols.append({
                    "name": name,
                    "type": node.type,
                    "start_line": node.start_point[0] + 1,
                    "end_line": node.end_point[0] + 1,
                    "language": lang_name,
                })
        # arrow function の検出 (const x = (...) => {...})
        if config.get("arrow_detect") and node.type == "variable_declarator":
            # 子に arrow_function があるか
            for child in node.children:
                if child.type == "arrow_function":
                    name_node = node.child_by_field_name("name")
                    if name_node:
                        symbols.append({
                            "name": name_node.text.decode("utf-8", errors="replace"),
                            "type": "arrow_function",
                            "start_line": node.start_point[0] + 1,
                            "end_line": child.end_point[0] + 1,
                            "language": lang_name,
                        })
                    break

        # import文
        if node.type in config.get("references", []):
            text = node.text.decode("utf-8", errors="replace").strip()
            symbols.append({
                "name": text,
                "type": node.type,
                "start_line": node.start_point[0] + 1,
                "end_line": node.end_point[0] + 1,
                "language": lang_name,
                "is_import": True,
            })

        for child in node.children:
            _walk_for_definitions(child, depth + 1)

    _walk_for_definitions(tree.root_node)
    return symbols

# ========================================================
# MCPサーバー初期化
# ========================================================

mcp = FastMCP("unified-mcp")

# ========================================================
# ChromaDB初期化（1つのDBで全プロジェクト管理・遅延ロード）
# ========================================================
# 重い依存（chromadb / sentence-transformers / torch）は最初のRAG呼び出しまで
# 読み込まない。RAG無効時やライブラリ未インストール時でもサーバーは起動する。

chroma_client = None
embedding_func = None
collection = None


class RAGUnavailable(RuntimeError):
    """RAGが無効、または依存ライブラリが未インストールのときに送出。"""


def _get_collection():
    """ChromaDBのcollectionを遅延生成して返す。無効/未導入なら RAGUnavailable。"""
    global chroma_client, embedding_func, collection
    if not ENABLE_RAG:
        raise RAGUnavailable(
            "RAG は無効化されています（ENABLE_RAG=0）。意味検索を使うには ENABLE_RAG=1 で起動してください。"
        )
    if collection is not None:
        return collection
    try:
        import chromadb
        from chromadb.utils import embedding_functions
    except ImportError as e:
        raise RAGUnavailable(
            "RAG 用ライブラリが未インストールです。`pip install chromadb sentence-transformers` を実行してください。"
            f"（詳細: {e}）"
        )
    os.makedirs(DB_DIR, exist_ok=True)
    chroma_client = chromadb.PersistentClient(path=DB_DIR)
    embedding_func = embedding_functions.SentenceTransformerEmbeddingFunction(
        model_name=EMBED_MODEL
    )
    collection = chroma_client.get_or_create_collection(
        name="unified_docs",
        embedding_function=embedding_func,
    )
    return collection

# ========================================================
# 更新状態管理（スレッドセーフ）
# ========================================================

_update_lock = threading.Lock()
update_status = {
    "is_running": False,
    "progress": "",
    "last_result": ""
}


# ========================================================
# ユーティリティ
# ========================================================

def _is_excluded(path: Path) -> bool:
    return any(part in EXCLUDE_DIRS for part in path.parts)

# .env系ファイルを弾くパターン
_ENV_FILE_PATTERN = re.compile(r'^\.env(\..+)?$')

def _is_sensitive_file(path: Path) -> bool:
    """シークレット系ファイル（.env等）をLLMに読ませない。"""
    return bool(_ENV_FILE_PATTERN.match(path.name))

def _calculate_file_md5(filepath: Path) -> str:
    md5 = hashlib.md5()
    with filepath.open("rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            md5.update(chunk)
    return md5.hexdigest()


def _get_projects(root: str) -> list[str]:
    """Get project directory names under the specified root."""
    base = Path(root)
    if not base.exists():
        return []
    return sorted([d.name for d in base.iterdir() if d.is_dir() and not _is_excluded(d)])


def _resolve_apps(project: str, filepath: str = "") -> tuple[Path, Path] | str:
    """Resolve an apps path and return (root, target)."""
    base = Path(APPS_ROOT) / project
    if not base.exists():
        available = _get_projects(APPS_ROOT)
        return f"❌ project '{project}' not found in apps. available: {available}"
    root = base.resolve()
    target = (root / filepath).resolve() if filepath else root
    try:
        target.relative_to(root)
    except ValueError:
        return "❌ target path is outside apps project root."
    return root, target


def _resolve_docs(project: str, filepath: str = "") -> tuple[Path, Path] | str:
    """Resolve a docs path and return (root, target)."""
    base = Path(DOCS_ROOT) / project
    if not base.exists():
        available = _get_projects(DOCS_ROOT)
        return f"❌ docs project '{project}' not found. available: {available}"
    root = base.resolve()
    target = (root / filepath).resolve() if filepath else root
    try:
        target.relative_to(root)
    except ValueError:
        return "❌ target path is outside docs project root."
    return root, target


def _is_searchable_text_file(filepath: Path) -> bool:
    """[FIX⑥] テキスト検索対象かどうかを判定する"""
    if filepath.suffix.lower() in TEXT_SEARCH_EXTENSIONS:
        return True
    if filepath.name in TEXT_SEARCH_NAMES:
        return True
    return False


def _truncate_content(text: str, max_chars: int = MAX_READ_CHARS) -> tuple[str, bool]:
    """[FIX①] 長すぎるコンテンツを先頭+末尾に切り詰める"""
    if len(text) <= max_chars:
        return text, False

    head_size = int(max_chars * 0.7)
    tail_size = max_chars - head_size
    truncated = (
        text[:head_size]
        + f"\n\n... [省略: 全{len(text)}文字中、中間部分を省略しています。"
        + f" read_pages() やファイルパス指定で部分取得してください] ...\n\n"
        + text[-tail_size:]
    )
    return truncated, True


# ========================================================
# テキスト抽出（多拡張子対応）
# ========================================================

def _read_text_file(filepath: str) -> str:
    """Read a text file with encoding fallback."""
    for enc in ("utf-8", "cp932", "latin-1"):
        try:
            with open(filepath, "r", encoding=enc) as f:
                return f.read()
        except (UnicodeDecodeError, LookupError):
            continue
    with open(filepath, "rb") as f:
        return f.read().decode("utf-8", errors="replace")


def _extract_pdf(filepath: str) -> str:
    if not _PYMUPDF_AVAILABLE:
        raise RuntimeError("PyMuPDF (fitz) が未インストールのため PDF を読めません。`pip install pymupdf` を実行してください。")
    doc = fitz.open(filepath)
    text = ""
    for page_num, page in enumerate(doc, 1):
        text += f"\n[PDF_PAGE_{page_num}]\n"
        text += page.get_text()
    doc.close()
    return text


def _extract_pptx(filepath: str) -> str:
    from pptx import Presentation
    prs = Presentation(filepath)
    text = ""
    for i, slide in enumerate(prs.slides, 1):
        text += f"\n--- スライド {i} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def _extract_docx(filepath: str) -> str:
    from docx import Document
    doc = Document(filepath)
    lines = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            lines.append("\t".join(c.text for c in row.cells))
    return "\n".join(lines)


def _extract_xlsx(filepath: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(filepath, data_only=True)
    text = ""
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text += f"\n--- シート: {sheet_name} ---\n"
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join([str(c) if c is not None else "" for c in row])
            if row_text.strip():
                text += row_text + "\n"
    return text


def _extract_xls(filepath: str) -> str:
    import xlrd
    wb = xlrd.open_workbook(filepath)
    text = ""
    for sheet in wb.sheets():
        text += f"\n--- シート: {sheet.name} ---\n"
        for row_idx in range(sheet.nrows):
            row_text = "\t".join([str(cell.value) for cell in sheet.row(row_idx)])
            if row_text.strip():
                text += row_text + "\n"
    return text


def extract_text(filepath: str) -> tuple[str, str | None]:
    """Extract text from file. Supports binary office/pdf formats too."""
    ext = Path(filepath).suffix.lower()
    binary_extractors = {
        ".pdf":  _extract_pdf,
        ".pptx": _extract_pptx,
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".xls":  _extract_xls,
    }
    try:
        if ext in binary_extractors:
            return binary_extractors[ext](filepath), None
        else:
            return _read_text_file(filepath), None
    except Exception as e:
        return "", f"[{ext}読み込みエラー] {Path(filepath).name}: {e}"


# ========================================================
# チャンク化
# ========================================================

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[dict]:
    page_markers = list(re.finditer(r'\[PDF_PAGE_(\d+)\]\n?', text))
    position_map = []

    if page_markers:
        clean_pos = 0
        original_pos = 0
        for match in page_markers:
            marker_start = match.start()
            page_num = int(match.group(1))
            text_before = text[original_pos:marker_start]
            clean_pos += len(text_before)
            position_map.append((clean_pos, match.end(), page_num))
            original_pos = match.end()
        if original_pos < len(text) and page_markers:
            final_page = int(page_markers[-1].group(1))
            position_map.append((clean_pos + len(text[original_pos:]), len(text), final_page))

    def get_page_range(start_pos, end_pos):
        if not position_map:
            return None, None
        pages = {pn for cp, _, pn in position_map if start_pos <= cp <= end_pos or cp <= start_pos}
        return (min(pages), max(pages)) if pages else (None, None)

    clean_text = re.sub(r'\[PDF_PAGE_\d+\]\n?', '', text)
    sentences = re.split(r'(?<=[。\n])|(?<=\.\s)', clean_text)
    sentences = [s for s in sentences if s.strip()]

    chunks: list[dict] = []
    current_chunk = ""
    current_start = 0

    for sentence in sentences:
        if len(current_chunk) + len(sentence) <= chunk_size:
            current_chunk += sentence
        else:
            if current_chunk.strip():
                end = current_start + len(current_chunk)
                ps, pe = get_page_range(current_start, end)
                chunks.append({"text": current_chunk.strip(), "page_start": ps, "page_end": pe})
                current_start = end
            if len(sentence) > chunk_size:
                start = 0
                while start < len(sentence):
                    piece = sentence[start:start + chunk_size]
                    if piece.strip():
                        ps, pe = get_page_range(current_start + start, current_start + start + len(piece))
                        chunks.append({"text": piece.strip(), "page_start": ps, "page_end": pe})
                    start += (chunk_size - overlap)
                current_start += len(sentence)
                current_chunk = ""
            else:
                overlap_text = chunks[-1]["text"][-overlap:] if chunks else ""
                current_chunk = overlap_text + sentence
                current_start = current_start - len(overlap_text)

    if current_chunk.strip():
        end = current_start + len(current_chunk)
        ps, pe = get_page_range(current_start, end)
        chunks.append({"text": current_chunk.strip(), "page_start": ps, "page_end": pe})

    return chunks


# ========================================================
# RAGインデックス作成
# ========================================================

def _scan_and_index(project: str | None = None) -> str:
    """docs配下を差分スキャンしてChromaDBを更新する。project指定で対象を絞る。"""
    try:
        col = _get_collection()
    except RAGUnavailable as e:
        return f"ERROR: {e}"

    docs_base = Path(DOCS_ROOT)
    if not docs_base.exists():
        return f"ERROR: docs folder not found: {DOCS_ROOT}"

    target_projects = [project] if project else _get_projects(DOCS_ROOT)
    if not target_projects:
        return "WARN: no projects found under docs."

    with _update_lock:
        update_status["progress"] = "既存データを確認中..."

    total_success, total_errors, total_skipped = 0, 0, 0
    error_msgs = []

    for proj in target_projects:
        proj_path = docs_base / proj
        if not proj_path.exists():
            continue

        existing = col.get(where={"project": proj}, include=["metadatas"])
        existing_by_path: dict[str, dict] = {}
        for id_, meta in zip(existing.get("ids", []), existing.get("metadatas", [])):
            if not meta:
                continue
            full_path = meta.get("full_path")
            if not full_path:
                continue
            if full_path not in existing_by_path:
                existing_by_path[full_path] = {
                    "ids": [],
                    "content_hash": meta.get("content_hash", ""),
                }
            existing_by_path[full_path]["ids"].append(id_)

        # docs/配下 + EXTRA_RAG_PATHS のファイルをまとめて処理
        files_to_index: list[tuple[Path, str, str]] = []  # (filepath, rel_path_str, folder)

        for filepath in proj_path.rglob("*"):
            if filepath.is_file() and not _is_excluded(filepath):
                rel = filepath.relative_to(proj_path)
                rel_str = str(rel).replace("\\", "/")
                folder = str(rel.parent) if rel.parent != Path(".") else ""
                files_to_index.append((filepath, rel_str, folder))

        for extra_dir_str in EXTRA_RAG_PATHS.get(proj, []):
            extra_dir = Path(extra_dir_str)
            if not extra_dir.exists():
                continue
            for filepath in extra_dir.rglob("*"):
                if filepath.is_file() and not _is_excluded(filepath):
                    files_to_index.append((filepath, "extra/" + filepath.name, "extra"))

        seen_paths: set[str] = set()

        for i, (filepath, rel_path_str, folder) in enumerate(files_to_index, 1):
            with _update_lock:
                update_status["progress"] = f"[{proj}] {i}/{len(files_to_index)} - {filepath.name}"

            seen_paths.add(rel_path_str)

            file_md5 = _calculate_file_md5(filepath)
            existing_entry = existing_by_path.get(rel_path_str)
            if existing_entry and existing_entry.get("content_hash") == file_md5:
                total_skipped += 1
                continue

            if existing_entry and existing_entry["ids"]:
                col.delete(ids=existing_entry["ids"])

            text, error = extract_text(str(filepath))
            if error or not text.strip():
                error_msgs.append(error or f"{filepath.name}: 抽出失敗")
                total_errors += 1
                continue

            file_path_hash = hashlib.md5(f"{proj}/{rel_path_str}".encode()).hexdigest()[:8]
            file_content_hash = file_md5[:8]

            # --- [FIX④] バッチ登録に変更 ---
            ext = filepath.suffix.lower()
            c_size = PDF_CHUNK_SIZE if ext == ".pdf" else CHUNK_SIZE
            c_overlap = PDF_CHUNK_OVERLAP if ext == ".pdf" else CHUNK_OVERLAP
            chunks = chunk_text(text, chunk_size=c_size, overlap=c_overlap)
            batch_ids = []
            batch_docs = []
            batch_metas = []

            for idx, chunk_data in enumerate(chunks):
                metadata = {
                    "project": proj,
                    "filename": filepath.name,
                    "folder": folder,
                    "full_path": rel_path_str,
                    "content_hash": file_md5,
                    "chunk_index": idx,
                }
                if chunk_data["page_start"] is not None:
                    metadata["page_start"] = chunk_data["page_start"]
                    metadata["page_end"] = chunk_data["page_end"]

                batch_ids.append(f"{proj}_{file_path_hash}_{file_content_hash}_{idx}")
                batch_docs.append(f"passage: {chunk_data['text']}")
                batch_metas.append(metadata)

            # ChromaDB のバッチ上限は約5461件なので安全に分割
            BATCH_LIMIT = 5000
            for start in range(0, len(batch_ids), BATCH_LIMIT):
                end = start + BATCH_LIMIT
                col.add(
                    ids=batch_ids[start:end],
                    documents=batch_docs[start:end],
                    metadatas=batch_metas[start:end],
                )

            total_success += 1

        deleted_paths = set(existing_by_path.keys()) - seen_paths
        ids_to_delete = []
        for removed_path in deleted_paths:
            ids_to_delete.extend(existing_by_path[removed_path]["ids"])
        if ids_to_delete:
            col.delete(ids=ids_to_delete)

    result = (
        f"Index refresh completed\n"
        f"  target projects: {target_projects}\n"
        f"  updated: {total_success} files / skipped: {total_skipped} files / errors: {total_errors} files\n"
        f"  timestamp: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    )
    if error_msgs:
        result += "\n\n[error details]\n" + "\n".join(f"  - {m}" for m in error_msgs[:5])
        if len(error_msgs) > 5:
            result += f"\n  ... 他 {len(error_msgs) - 5} 件"

    return result


def _do_refresh(project: str | None):
    try:
        result = _scan_and_index(project)
    except Exception as e:
        result = f"❌ エラー: {e}"
    finally:
        with _update_lock:
            update_status["last_result"] = result
            update_status["is_running"] = False
            update_status["progress"] = ""


# ========================================================
# MCPツール: apps側（読み書き）
# ========================================================

@mcp.tool()
def list_files(project: str, subpath: str = "") -> str:
    """appsプロジェクト内のファイル一覧を表示する。
    サブパス未指定時はルート直下の1階層のみ表示。
    サブパスを指定するとその直下1階層を表示。
    ツリー全体が必要な場合は subpath="..." は使わず繰り返し呼び出すこと。
    Args:
        project: プロジェクト名 (例: "research", "akatsuki", "teach")
                 利用可能なプロジェクトは list_projects() で確認できる
        subpath: サブフォルダの相対パス（省略時はルート表示）
    """
    result = _resolve_apps(project, subpath)
    if isinstance(result, str):
        return result
    root, target = result

    if not target.exists():
        return f"❌ 存在しません: {subpath or '(root)'}"
    if not target.is_dir():
        return f"❌ ディレクトリではありません: {subpath}"

    items = []
    for item in sorted(target.iterdir()):
        if _is_excluded(item):
            continue
        rel = item.relative_to(root)
        kind = "📁" if item.is_dir() else "📄"
        size = f" ({round(item.stat().st_size / 1024, 1)}KB)" if item.is_file() else ""
        items.append(f"{kind} {rel}{size}")

    header = f"📁 [{project}] {target}"
    if not items:
        return f"{header}\n(empty)"
    return header + "\n" + "\n".join(items)


@mcp.tool()
def read_file(project: str, filepath: str) -> str:
    """Read a file under apps project and return JSON payload.
    [FIX①] 長大なファイルは自動的に先頭+末尾に切り詰めて返す。
    """
    result = _resolve_apps(project, filepath)
    if isinstance(result, str):
        return result
    root, target = result

    if not target.exists():
        return f"❌ file not found: {filepath}"
    if target.is_dir():
        return f"❌ path is a directory: {filepath}. Use list_files() first."
    if _is_excluded(target):
        return "❌ this path is excluded."
    if _is_sensitive_file(target):          # ← 追加
        return "❌ .env系ファイルは読み取り禁止です。"

    text, error = extract_text(str(target))
    if error:
        return f"❌ {error}"

    text, was_truncated = _truncate_content(text)

    payload = {
        "project":  project,
        "filepath": filepath,
        "filename": target.name,
        "size_kb":  round(target.stat().st_size / 1024, 1),
        "content":  text,
    }
    if was_truncated:
        payload["truncated"] = True
        payload["total_chars"] = len(text)
    return json.dumps(payload, ensure_ascii=False, indent=2)


@mcp.tool()
def write_file(project: str, filepath: str, content: str) -> str:
    """appsプロジェクト内のファイルを書き込む（新規作成・上書き）。バックアップを自動作成。
    ※ 大きなファイルの部分編集には patch_file() を使うこと。
    Args:
        project:  プロジェクト名 (例: "research", "akatsuki", "teach")
        filepath: プロジェクトルートからの相対パス (例: "src/pages/index.tsx", "main.py")
        content:  書き込み内容（全文）
    """
    result = _resolve_apps(project, filepath)
    if isinstance(result, str):
        return result
    root, target = result

    if _is_excluded(target):
        return "❌ this path is excluded."

    target.parent.mkdir(parents=True, exist_ok=True)

    backup_msg = ""
    if target.exists():
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup_root = root / ".backups"
        backup_rel = target.relative_to(root)
        backup = backup_root / backup_rel.parent / f"{backup_rel.name}.{ts}.bak"
        try:
            backup.parent.mkdir(parents=True, exist_ok=True)
            backup.write_bytes(target.read_bytes())
            backup_msg = f"\nBackup: {backup.relative_to(root)}"
        except Exception as e:
            backup_msg = f"\nBackup failed (write continues): {e}"

    try:
        target.write_text(content, encoding="utf-8")
        return f"✅ saved [{project}] {filepath}{backup_msg}"
    except Exception as e:
        return f"❌ バックアップ失敗のため書き込みを中止: {e}\nBackup path: {backup}"


# --- [FIX⑤+NEW②] 差分編集ツール（anchor-based + line-based） ---
@mcp.tool()
def patch_file(project: str, filepath: str, patches: str) -> str:
    """appsプロジェクト内のファイルを差分編集する。全文書き換え不要。バックアップを自動作成。

    write_file() と違い、変更箇所だけを指定できるため、
    LLMの出力トークンを節約でき、大きなファイルの編集に適している。

    ■ 推奨: anchor（文字列検索）方式 — 行番号ズレに強い
      {"action": "anchor_replace", "search": "検索文字列", "replace": "置換文字列"}
      {"action": "anchor_insert", "search": "この文字列の直後に挿入", "content": "挿入内容"}
      {"action": "anchor_delete", "search": "削除する文字列"}
      ※ search はファイル内で一意にマッチする必要がある（0件/複数件はエラー）。
      ※ search は複数行にまたがってもOK（改行を含む文字列で検索可能）。

    ■ フォールバック: line（行番号）方式
      {"action": "replace", "start_line": 10, "end_line": 15, "content": "新しい内容"}
      {"action": "insert", "after_line": 5, "content": "挿入する内容"}
      {"action": "delete", "start_line": 10, "end_line": 12}
      ※ 行番号は1始まり。end_lineはその行を含む。
      ※ 複数パッチは末尾行から適用すること（行番号ズレ防止）。

    Args:
        project:  プロジェクト名
        filepath: プロジェクトルートからの相対パス
        patches:  JSON配列文字列。上記いずれかの形式。混在も可。
    """
    result = _resolve_apps(project, filepath)
    if isinstance(result, str):
        return result
    root, target = result

    if not target.exists():
        return f"❌ file not found: {filepath}"
    if _is_excluded(target):
        return "❌ this path is excluded."

    try:
        patch_list = json.loads(patches)
        if not isinstance(patch_list, list):
            return "❌ patches must be a JSON array."
    except json.JSONDecodeError as e:
        return f"❌ invalid JSON in patches: {e}"

    # 既存内容を読み込み
    text, error = extract_text(str(target))
    if error:
        return f"❌ {error}"

    # バックアップ
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    backup_root = root / ".backups"
    backup_rel = target.relative_to(root)
    backup = backup_root / backup_rel.parent / f"{backup_rel.name}.{ts}.bak"
    try:
        backup.parent.mkdir(parents=True, exist_ok=True)
        backup.write_bytes(target.read_bytes())
    except Exception as e:
        return f"❌ バックアップ失敗のため編集を中止: {e}\nBackup path: {backup}"

    # anchor系とline系を分離
    anchor_patches = []
    line_patches = []
    errors = []
    for i, p in enumerate(patch_list):
        action = p.get("action", "")
        if action.startswith("anchor_"):
            anchor_patches.append((i, p))
        elif action in ("replace", "insert", "delete"):
            line_patches.append((i, p))
        else:
            errors.append(f"patch[{i}]: unknown action '{action}'")

    applied = 0

    # === anchor方式を先に適用（全文文字列ベースなので順序無関係） ===
    for i, patch in anchor_patches:
        action = patch["action"]
        search = patch.get("search", "")
        if not search:
            errors.append(f"patch[{i}]: 'search' is required for anchor actions")
            continue

        count = text.count(search)
        if count == 0:
            errors.append(f"patch[{i}]: search string not found (0 matches)")
            continue
        if count > 1:
            errors.append(
                f"patch[{i}]: search string matched {count} times (must be unique). "
                f"Include more surrounding context in 'search' to narrow it down."
            )
            continue

        try:
            if action == "anchor_replace":
                replacement = patch.get("replace", "")
                text = text.replace(search, replacement, 1)
                applied += 1
            elif action == "anchor_insert":
                content = patch.get("content", "")
                insert_pos = text.index(search) + len(search)
                text = text[:insert_pos] + content + text[insert_pos:]
                applied += 1
            elif action == "anchor_delete":
                text = text.replace(search, "", 1)
                applied += 1
            else:
                errors.append(f"patch[{i}]: unknown anchor action '{action}'")
        except Exception as ex:
            errors.append(f"patch[{i}]: {ex}")

    # === line方式を適用（anchor適用後のtextに対して） ===
    if line_patches:
        lines = text.splitlines(keepends=True)

        # 末尾行から適用（行番号ズレ防止）
        def sort_key(item):
            _, p = item
            return p.get("start_line", p.get("after_line", 0))
        line_patches.sort(key=sort_key, reverse=True)

        for i, patch in line_patches:
            action = patch.get("action", "")
            try:
                if action == "replace":
                    s = patch["start_line"] - 1
                    e = patch["end_line"]
                    new_content = patch["content"]
                    if not new_content.endswith("\n"):
                        new_content += "\n"
                    lines[s:e] = [new_content]
                    applied += 1
                elif action == "insert":
                    after = patch["after_line"]
                    new_content = patch["content"]
                    if not new_content.endswith("\n"):
                        new_content += "\n"
                    lines.insert(after, new_content)
                    applied += 1
                elif action == "delete":
                    s = patch["start_line"] - 1
                    e = patch["end_line"]
                    del lines[s:e]
                    applied += 1
            except (KeyError, IndexError) as ex:
                errors.append(f"patch[{i}]: {ex}")

        text = "".join(lines)

    # 書き込み
    try:
        target.write_text(text, encoding="utf-8")
    except Exception as e:
        return f"❌ write error: {e}"

    msg = f"✅ patched [{project}] {filepath} ({applied} patches applied)"
    if errors:
        msg += "\n⚠️ errors:\n" + "\n".join(f"  - {e}" for e in errors)
    return msg


# --- [FIX②] search_in_files を高速化 ---
@mcp.tool()
def search_in_files(project: str, query: str) -> str:
    """appsプロジェクト内のテキストファイルからキーワードをgrep検索する。
    Args:
        project: プロジェクト名 (例: "research", "akatsuki", "teach")
        query:   検索したいキーワード（関数名・変数名・クラス名など）
    """
    result = _resolve_apps(project, "")
    if isinstance(result, str):
        return result
    root, _ = result

    # [FIX⑥] ホワイトリスト＋サイズフィルタで対象を絞る
    target_files = []
    for filepath in sorted(root.rglob("*")):
        if not filepath.is_file() or _is_excluded(filepath):
            continue
        if filepath.suffix.lower() in BINARY_EXTENSIONS:
            continue
        if not _is_searchable_text_file(filepath):
            continue
        if _is_sensitive_file(filepath):    # ← 追加
            continue
        try:
            if filepath.stat().st_size > MAX_SEARCH_FILE_SIZE:
                continue
        except OSError:
            continue
        target_files.append(filepath)

    query_lower = query.lower()
    results = []

    # [FIX②] 並列ファイル読み込みで高速化
    def _search_one(fp: Path):
        try:
            text = _read_text_file(str(fp))
            matches = [
                {"line": i + 1, "text": line.rstrip()[:200]}  # 行の長さも制限
                for i, line in enumerate(text.splitlines())
                if query_lower in line.lower()
            ]
            if matches:
                return {
                    "file":          str(fp.relative_to(root)),
                    "matches":       matches[:5],
                    "total_matches": len(matches),
                }
        except Exception:
            pass
        return None

    with ThreadPoolExecutor(max_workers=min(8, os.cpu_count() or 4)) as executor:
        futures = {executor.submit(_search_one, fp): fp for fp in target_files}
        for future in as_completed(futures):
            hit = future.result()
            if hit:
                results.append(hit)

    # ファイル名でソート
    results.sort(key=lambda r: r["file"])

    if not results:
        return f"No matches for '{query}' in [{project}]. (searched {len(target_files)} files)"
    return json.dumps(results, ensure_ascii=False, indent=2)


# ========================================================
# MCPツール: docs側（RAG検索）
# ========================================================

@mcp.tool()
def search_docs(query: str, project: str = "") -> str:
    """docs内の文書からキーワードや質問で意味検索する。
    💡 結果はスニペット。詳細は full_path を使って read_document() を呼ぶこと。
    Args:
        query:   検索したいキーワードや質問
        project: プロジェクト名で絞り込み（省略で全プロジェクト横断検索）
                 例: "research", "akatsuki"
    """
    try:
        col = _get_collection()
    except RAGUnavailable as e:
        return f"❌ {e}"

    # --- [FIX③] project指定時はChromaDBのwhereフィルタを使う ---
    query_kwargs: dict = {
        "query_texts": [f"query: {query}"],
        "n_results": SEARCH_N_RESULTS,
        "include": ["documents", "metadatas", "distances"],
    }
    if project:
        query_kwargs["where"] = {"project": project}

    results = col.query(**query_kwargs)
    if not results["ids"][0]:
        return json.dumps([], ensure_ascii=False)

    seen: dict[str, dict] = {}
    for doc, meta, dist in zip(
        results["documents"][0],
        results["metadatas"][0],
        results["distances"][0],
    ):
        if dist > SEARCH_DISTANCE_THRESHOLD:
            continue

        fp = meta.get("full_path", meta.get("filename", "不明"))
        proj = meta.get("project", "")
        key = f"{proj}/{fp}"
        if key in seen:
            continue

        snippet = doc.removeprefix("passage: ").strip().replace("\n", " ")
        if len(snippet) > 150:
            snippet = snippet[:150] + "…"

        info = {
            "project":   proj,
            "filename":  meta.get("filename", "不明"),
            "folder":    meta.get("folder", ""),
            "full_path": fp,
            "snippet":   snippet,
            "score":     round(1 - dist, 3),
        }
        if "page_start" in meta:
            ps, pe = meta["page_start"], meta["page_end"]
            info["page_info"] = f"p.{ps}" if ps == pe else f"p.{ps}-{pe}"
        seen[key] = info
        if len(seen) >= SEARCH_MAX_FILES:
            break

    return json.dumps(list(seen.values()), ensure_ascii=False, indent=2)


@mcp.tool()
def read_document(project: str, full_path: str) -> str:
    """docs内の文書の全文を取得する。
    [FIX①] 長大な文書は自動的に切り詰めて返す。全文が必要ならread_pages()でページ範囲指定。
    Args:
        project:   プロジェクト名 (例: "research", "akatsuki")
        full_path: search_docs の結果に含まれる full_path の値
    """
    result = _resolve_docs(project, full_path)
    if isinstance(result, str):
        return result
    _, target = result

    normalized = unicodedata.normalize("NFD", str(target))
    target = Path(normalized)

    if not target.exists():
        return f"❌ ファイルが見つかりません: {full_path}"

    text, error = extract_text(str(target))
    if error:
        return f"❌ {error}"
    if not text.strip():
        return f"⚠ テキストを抽出できませんでした: {full_path}"

    text, was_truncated = _truncate_content(text)

    payload = {
        "project":   project,
        "filename":  target.name,
        "full_path": full_path,
        "content":   text,
    }
    if was_truncated:
        payload["truncated"] = True
        payload["note"] = "文書が長いため切り詰めました。read_pages()でページ範囲を指定して取得できます。"
    return json.dumps(payload, ensure_ascii=False, indent=2)


@mcp.tool()
def read_pages(project: str, full_path: str, start_page: int, end_page: int) -> str:
    """PDFの指定ページ範囲だけを取得する。
    Args:
        project:    プロジェクト名
        full_path:  search_docs の full_path
        start_page: 開始ページ（1始まり）
        end_page:   終了ページ（start_pageと同値で1ページのみ）
    """
    result = _resolve_docs(project, full_path)
    if isinstance(result, str):
        return result
    _, target = result

    if not target.exists():
        return f"❌ file not found: {full_path}"
    if target.suffix.lower() != ".pdf":
        return "❌ target is not a PDF. Use read_document()."
    if not _PYMUPDF_AVAILABLE:
        return "❌ PyMuPDF (fitz) が未インストールです。`pip install pymupdf` を実行してください。"

    try:
        doc = fitz.open(str(target))
        total = len(doc)
        if not (1 <= start_page <= total) or not (start_page <= end_page <= total):
            doc.close()
            return f"❌ invalid page range. total pages: {total}"
        parts = []
        for pn in range(start_page - 1, end_page):
            parts.append(f"\n━━━ Page {pn+1} / {total} ━━━\n")
            parts.append(doc[pn].get_text())
        doc.close()
        return json.dumps({
            "project":     project,
            "filename":    target.name,
            "full_path":   full_path,
            "page_range":  f"{start_page}-{end_page}",
            "total_pages": total,
            "content":     "".join(parts),
        }, ensure_ascii=False, indent=2)
    except Exception as e:
        return f"❌ 読み込みエラー: {e}"


@mcp.tool()
def list_docs(project: str = "") -> str:
    """RAGに読み込まれている文書の一覧を表示する。
    Args:
        project: プロジェクト名で絞り込み（省略で全プロジェクト表示）
    """
    try:
        col = _get_collection()
    except RAGUnavailable as e:
        return f"❌ {e}"

    # --- [FIX②] project指定時はwhereフィルタを使い、全取得を避ける ---
    get_kwargs: dict = {"include": ["metadatas"]}
    if project:
        get_kwargs["where"] = {"project": project}
    existing = col.get(**get_kwargs)

    if not existing["ids"]:
        msg = f"No indexed documents for project '{project}'." if project else "No indexed documents yet."
        return msg + " Run refresh_database()."

    files_by_proj: dict[str, dict[str, set]] = {}
    for meta in existing["metadatas"]:
        proj = meta.get("project", "不明")
        folder = meta.get("folder", "")
        fname  = meta.get("filename", "不明")
        files_by_proj.setdefault(proj, {}).setdefault(folder, set()).add(fname)

    output = "【docs 文書一覧】\n"
    total = 0
    for proj in sorted(files_by_proj):
        output += f"\n📂 {proj}\n"
        for folder in sorted(files_by_proj[proj]):
            label = f"  📁 {folder}/" if folder else "  📁 (ルート)"
            output += f"{label}\n"
            for f in sorted(files_by_proj[proj][folder]):
                output += f"    - {f}\n"
                total += 1
    output += f"\n合計: {total}ファイル"
    return output


@mcp.tool()
def refresh_database(project: str = "") -> str:
    """Refresh RAG index in the background."""
    with _update_lock:
        if update_status["is_running"]:
            return f"refresh already running: {update_status['progress']}\nUse check_update_status()."
        update_status.update({"is_running": True, "progress": "開始中...", "last_result": ""})

    target_proj = project if project else None
    threading.Thread(target=_do_refresh, args=(target_proj,), daemon=True).start()

    scope = f"[{project}]" if project else "[all projects]"
    return f"started refresh for {scope}.\nUse check_update_status()."


@mcp.tool()
def check_update_status() -> str:
    """Show current refresh progress and last result."""
    with _update_lock:
        is_running = update_status["is_running"]
        progress   = update_status["progress"]
        last       = update_status["last_result"]

    if is_running:
        return f"⏳ 更新中...\n現在: {progress}"
    elif last:
        return f"最新の更新結果:\n{last}"
    return "No refresh has been run yet. Call refresh_database()."


# ========================================================
# MCPツール: [NEW①③] Symbol解析 + 範囲読み
# ========================================================

@mcp.tool()
def index_symbols(project: str) -> str:
    """appsプロジェクト内の全ファイルからシンボル（関数・クラス・interface等）を抽出して一覧を返す。

    tree-sitter でAST解析を行うため、文字列検索では拾えない正確な定義位置がわかる。
    対応言語: Python (.py), TypeScript (.ts/.tsx), JavaScript (.js/.jsx), MATLAB (.m)
    ※ 未対応拡張子のファイルはスキップされる。

    Args:
        project: プロジェクト名
    """
    result = _resolve_apps(project, "")
    if isinstance(result, str):
        return result
    root, _ = result

    err = _require_tree_sitter()
    if err:
        return err
    all_symbols: dict[str, list] = {}  # filepath -> [symbols]
    file_count = 0
    skip_count = 0

    for filepath in sorted(root.rglob("*")):
        if not filepath.is_file() or _is_excluded(filepath):
            continue
        symbols = _extract_symbols_from_file(filepath)
        if symbols is None:
            skip_count += 1
            continue
        file_count += 1
        rel = str(filepath.relative_to(root)).replace("\\", "/")
        if symbols:
            all_symbols[rel] = symbols

    # 整形して返す
    output_entries = []
    for rel_path, syms in sorted(all_symbols.items()):
        defs = [s for s in syms if not s.get("is_import")]
        imports = [s for s in syms if s.get("is_import")]
        entry = {"file": rel_path, "definitions": [], "imports": []}
        for s in defs:
            entry["definitions"].append({
                "name": s["name"],
                "type": s["type"],
                "lines": f"{s['start_line']}-{s['end_line']}",
            })
        for s in imports:
            entry["imports"].append({
                "name": s["name"],
                "line": s["start_line"],
            })
        output_entries.append(entry)

    summary = {
        "project": project,
        "files_analyzed": file_count,
        "files_skipped": skip_count,
        "files_with_symbols": len(output_entries),
        "available_languages": sorted(_ext_to_lang.keys()),
        "symbols": output_entries,
    }
    return json.dumps(summary, ensure_ascii=False, indent=2)


@mcp.tool()
def find_symbol(project: str, name: str, kind: str = "") -> str:
    """プロジェクト内でシンボル（関数名・クラス名等）の定義箇所を検索する。

    search_in_files() と違い、AST解析に基づくため定義箇所を正確に特定できる。
    コメント内や文字列内の偶然の一致を除外できる。

    Args:
        project: プロジェクト名
        name:    検索するシンボル名（完全一致 or 部分一致）
        kind:    絞り込み: "function", "class", "interface", "arrow", "import" 等
                 省略で全種類検索
    """
    result = _resolve_apps(project, "")
    if isinstance(result, str):
        return result
    root, _ = result

    err = _require_tree_sitter()
    if err:
        return err
    name_lower = name.lower()
    matches = []

    for filepath in sorted(root.rglob("*")):
        if not filepath.is_file() or _is_excluded(filepath):
            continue
        symbols = _extract_symbols_from_file(filepath)
        if symbols is None:
            continue
        rel = str(filepath.relative_to(root)).replace("\\", "/")
        for s in symbols:
            sym_name = s["name"]
            # import文は名前全体で検索（"from pathlib import Path" の中に "Path" がある）
            if s.get("is_import"):
                if name_lower not in sym_name.lower():
                    continue
            else:
                if name_lower != sym_name.lower() and name_lower not in sym_name.lower():
                    continue
            if kind:
                if kind.lower() not in s["type"].lower():
                    continue
            matches.append({
                "file": rel,
                "name": sym_name,
                "type": s["type"],
                "lines": f"{s['start_line']}-{s['end_line']}",
                "start_line": s["start_line"],
                "end_line": s["end_line"],
            })

    if not matches:
        return f"No symbol '{name}' found in [{project}]."

    return json.dumps({
        "query": name,
        "kind_filter": kind or "(all)",
        "total_matches": len(matches),
        "matches": matches[:30],
    }, ensure_ascii=False, indent=2)


@mcp.tool()
def find_references(project: str, name: str) -> str:
    """プロジェクト内でシンボルが使われている（参照されている）箇所を検索する。

    find_symbol() が定義箇所を探すのに対し、こちらは呼び出し・参照箇所を探す。
    AST解析 + テキスト検索のハイブリッドで、import文やidentifierの出現箇所を返す。

    Args:
        project: プロジェクト名
        name:    検索するシンボル名（完全一致）
    """
    result = _resolve_apps(project, "")
    if isinstance(result, str):
        return result
    root, _ = result

    err = _require_tree_sitter()
    if err:
        return err
    references = []

    for filepath in sorted(root.rglob("*")):
        if not filepath.is_file() or _is_excluded(filepath):
            continue

        info = _get_parser_for_file(filepath)
        rel = str(filepath.relative_to(root)).replace("\\", "/")

        if info is not None:
            # AST解析で identifier ノードの出現を探す
            lang_name, parser, config = info
            try:
                source = filepath.read_bytes()
            except Exception:
                continue
            tree = parser.parse(source)
            source_text = source.decode("utf-8", errors="replace")

            file_refs = []

            def _find_identifiers(node):
                if node.type in ("identifier", "type_identifier", "property_identifier"):
                    node_text = node.text.decode("utf-8", errors="replace")
                    if node_text == name:
                        line_no = node.start_point[0] + 1
                        # コンテキスト（その行全体）
                        lines = source_text.splitlines()
                        if 0 <= node.start_point[0] < len(lines):
                            context = lines[node.start_point[0]].strip()[:200]
                        else:
                            context = ""
                        file_refs.append({
                            "line": line_no,
                            "context": context,
                        })
                for child in node.children:
                    _find_identifiers(child)

            _find_identifiers(tree.root_node)
            if file_refs:
                references.append({
                    "file": rel,
                    "language": lang_name,
                    "refs": file_refs[:10],
                    "total_refs": len(file_refs),
                })
        else:
            # AST非対応ファイルはテキスト検索にフォールバック
            if not _is_searchable_text_file(filepath):
                continue
            try:
                text = _read_text_file(str(filepath))
            except Exception:
                continue
            file_refs = []
            for i, line in enumerate(text.splitlines()):
                if name in line:
                    file_refs.append({
                        "line": i + 1,
                        "context": line.strip()[:200],
                    })
            if file_refs:
                references.append({
                    "file": rel,
                    "language": "text_search",
                    "refs": file_refs[:10],
                    "total_refs": len(file_refs),
                })

    if not references:
        return f"No references to '{name}' found in [{project}]."

    return json.dumps({
        "query": name,
        "files_with_refs": len(references),
        "references": references[:20],
    }, ensure_ascii=False, indent=2)


@mcp.tool()
def read_file_lines(project: str, filepath: str, start_line: int, end_line: int) -> str:
    """appsプロジェクト内のファイルの指定行範囲だけを取得する。

    read_file() の全文取得と違い、必要な範囲だけを読むためコンテキストを節約できる。
    find_symbol() で得た行番号と組み合わせて使うと効率的。

    Args:
        project:    プロジェクト名
        filepath:   プロジェクトルートからの相対パス
        start_line: 開始行（1始まり）
        end_line:   終了行（この行を含む。-1で末尾まで）
    """
    result = _resolve_apps(project, filepath)
    if isinstance(result, str):
        return result
    root, target = result

    if not target.exists():
        return f"❌ file not found: {filepath}"
    if target.is_dir():
        return f"❌ path is a directory: {filepath}"
    if _is_excluded(target):
        return "❌ this path is excluded."

    text, error = extract_text(str(target))
    if error:
        return f"❌ {error}"

    lines = text.splitlines(keepends=True)
    total = len(lines)

    if end_line == -1:
        end_line = total
    if not (1 <= start_line <= total):
        return f"❌ invalid start_line: {start_line}. total lines: {total}"
    if not (start_line <= end_line <= total):
        return f"❌ invalid end_line: {end_line}. total lines: {total}"

    selected = lines[start_line - 1 : end_line]

    return json.dumps({
        "project":    project,
        "filepath":   filepath,
        "start_line": start_line,
        "end_line":   end_line,
        "total_lines": total,
        "content":    "".join(selected),
    }, ensure_ascii=False, indent=2)


@mcp.tool()
def read_symbol(project: str, filepath: str, name: str) -> str:
    """ファイル内の特定シンボル（関数・クラス等）のソースコードだけを取得する。

    関数1つだけ読みたいとき、read_file() で全文を読む代わりにこれを使う。
    AST解析で定義の開始行〜終了行を特定し、その範囲だけ返す。

    Args:
        project:  プロジェクト名
        filepath: プロジェクトルートからの相対パス
        name:     読みたいシンボル名（関数名・クラス名等）
    """
    result = _resolve_apps(project, filepath)
    if isinstance(result, str):
        return result
    root, target = result

    if not target.exists():
        return f"❌ file not found: {filepath}"

    symbols = _extract_symbols_from_file(target)
    if symbols is None:
        return f"❌ {filepath} は symbol 解析に未対応の拡張子です。対応: {sorted(_ext_to_lang.keys())}"

    # 名前でマッチ（import以外）
    matched = [s for s in symbols if s["name"] == name and not s.get("is_import")]
    if not matched:
        # 部分一致でも試す
        matched = [s for s in symbols if name.lower() in s["name"].lower() and not s.get("is_import")]
    if not matched:
        available = [s["name"] for s in symbols if not s.get("is_import")]
        return f"❌ symbol '{name}' not found in {filepath}.\nAvailable: {available[:20]}"

    text, error = extract_text(str(target))
    if error:
        return f"❌ {error}"

    lines = text.splitlines(keepends=True)
    results = []
    for s in matched:
        start = s["start_line"] - 1
        end = s["end_line"]
        content = "".join(lines[start:end])
        results.append({
            "name": s["name"],
            "type": s["type"],
            "start_line": s["start_line"],
            "end_line": s["end_line"],
            "content": content,
        })

    return json.dumps({
        "project":  project,
        "filepath": filepath,
        "symbols":  results,
    }, ensure_ascii=False, indent=2)


# ========================================================
# MCPツール: プロジェクト管理
# ========================================================

@mcp.tool()
def list_projects() -> str:
    """List available projects in apps and docs."""
    apps_projects = _get_projects(APPS_ROOT)
    docs_projects = _get_projects(DOCS_ROOT)

    output = "Available projects\n\n"
    output += "apps (editable):\n"
    for p in apps_projects:
        tag = " *apps-only" if p in APPS_ONLY_PROJECTS else ""
        output += f"  - {p}{tag}\n"

    output += "\ndocs (RAG target):\n"
    for p in docs_projects:
        output += f"  - {p}\n"

    only_apps = set(apps_projects) - set(docs_projects) - APPS_ONLY_PROJECTS
    if only_apps:
        output += f"\n* apps only (not in docs): {sorted(only_apps)}"

    return output


# ========================================================
# MCPツール: 検証ループ（lint / typecheck / test / build）
# ========================================================

def _load_check_config(project_root: Path) -> dict | None:
    """プロジェクトルートの mcp_checks.json を読み込む。"""
    config_path = project_root / CHECK_CONFIG_FILENAME
    if not config_path.exists():
        return None
    try:
        return json.loads(config_path.read_text(encoding="utf-8"))
    except Exception:
        return None


@mcp.tool()
def list_checks(project: str) -> str:
    """プロジェクトで利用可能な検証コマンド一覧を表示する。

    プロジェクトルートに mcp_checks.json がある場合、その中の checks を表示する。
    ない場合は設定ファイルの作り方を案内する。

    Args:
        project: プロジェクト名
    """
    result = _resolve_apps(project, "")
    if isinstance(result, str):
        return result
    root, _ = result

    config = _load_check_config(root)
    if config is None:
        return (
            f"❌ [{project}] に {CHECK_CONFIG_FILENAME} がありません。\n\n"
            f"プロジェクトルートに以下の形式で作成してください:\n"
            f'{{\n'
            f'  "checks": {{\n'
            f'    "lint": "ruff check .",\n'
            f'    "typecheck": "pyright",\n'
            f'    "test": "pytest -x --tb=short",\n'
            f'    "build": "npm run build"\n'
            f'  }}\n'
            f'}}\n\n'
            f"キー名は自由です。コマンドはプロジェクトルートで実行されます。"
        )

    checks = config.get("checks", {})
    if not checks:
        return f"⚠️ [{project}] の {CHECK_CONFIG_FILENAME} に checks が定義されていません。"

    output = f"📋 [{project}] 利用可能な検証コマンド:\n"
    for name, cmd in checks.items():
        output += f"  - {name}: {cmd}\n"
    output += f"\n使い方: run_check(project=\"{project}\", check_name=\"<name>\")"
    return output


@mcp.tool()
def run_check(project: str, check_name: str) -> str:
    """プロジェクトの検証コマンドを実行して結果を返す。

    編集後の lint / typecheck / test / build 等を実行し、
    エラーがあれば出力を返す。LLMはこの結果を見て再修正できる。

    ※ セキュリティ: mcp_checks.json で事前定義されたコマンドのみ実行可能。
      任意コマンドは実行できない。

    Args:
        project:    プロジェクト名
        check_name: 実行する検証コマンド名（list_checks() で確認可能）
                    例: "lint", "test", "typecheck", "build"
    """
    result = _resolve_apps(project, "")
    if isinstance(result, str):
        return result
    root, _ = result

    config = _load_check_config(root)
    if config is None:
        return (
            f"❌ [{project}] に {CHECK_CONFIG_FILENAME} がありません。\n"
            f"list_checks(\"{project}\") で設定方法を確認してください。"
        )

    checks = config.get("checks", {})
    if check_name not in checks:
        available = list(checks.keys())
        return f"❌ check '{check_name}' は未定義です。利用可能: {available}"

    command = checks[check_name]
    timeout = config.get("timeout", CHECK_TIMEOUT_SECONDS)

    try:
        proc = subprocess.run(
            shlex.split(command),
            shell=False,
            cwd=str(root),
            capture_output=True,
            text=True,
            timeout=timeout,
            env={**os.environ, "FORCE_COLOR": "0", "NO_COLOR": "1"},
        )
    except subprocess.TimeoutExpired:
        return (
            f"⏰ [{project}] {check_name} がタイムアウトしました（{timeout}秒）。\n"
            f"コマンド: {command}"
        )
    except Exception as e:
        return f"❌ [{project}] {check_name} の実行に失敗: {e}"

    # 結果を整形
    status = "✅ PASS" if proc.returncode == 0 else f"❌ FAIL (exit code: {proc.returncode})"
    output = f"{status} [{project}] {check_name}\n"
    output += f"コマンド: {command}\n"

    stdout = proc.stdout.strip()
    stderr = proc.stderr.strip()

    # 出力が長すぎる場合は切り詰め
    MAX_OUTPUT = 5000
    if stdout:
        if len(stdout) > MAX_OUTPUT:
            stdout = stdout[:MAX_OUTPUT] + f"\n... (残り {len(stdout) - MAX_OUTPUT} 文字省略)"
        output += f"\n--- stdout ---\n{stdout}\n"
    if stderr:
        if len(stderr) > MAX_OUTPUT:
            stderr = stderr[:MAX_OUTPUT] + f"\n... (残り {len(stderr) - MAX_OUTPUT} 文字省略)"
        output += f"\n--- stderr ---\n{stderr}\n"

    if not stdout and not stderr:
        output += "\n(出力なし)"

    return output


# ========================================================
# 起動
# ========================================================

if __name__ == "__main__":
    # 作業フォルダが無ければ作る（初回起動でも tools が落ちないように）
    os.makedirs(APPS_ROOT, exist_ok=True)
    os.makedirs(DOCS_ROOT, exist_ok=True)

    print("unified-mcp サーバーを起動します...", file=sys.stderr)
    print(f"  WORKSPACE_ROOT = {MCP_ROOT}", file=sys.stderr)
    print(f"  apps           = {APPS_ROOT}", file=sys.stderr)
    print(f"  docs           = {DOCS_ROOT}", file=sys.stderr)
    print(f"  ENABLE_RAG     = {ENABLE_RAG}  (model: {EMBED_MODEL})", file=sys.stderr)
    print(f"  PyMuPDF        = {_PYMUPDF_AVAILABLE}", file=sys.stderr)
    mcp.run()
